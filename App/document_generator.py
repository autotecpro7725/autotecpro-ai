
"""
AutoTecPro AI document generator framework (v386 conversation export compatibility).

This module is independent of Streamlit, OpenAI, Supabase, WooCommerce, and
workspace UI code. It generates real downloadable files for PDF, DOCX, PPTX,
XLSX, and CSV and serializes them into chat-history-safe records.
"""

from __future__ import annotations

import base64
import csv
import io
import json
import re
import textwrap
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Iterable

DOCUMENT_MARKER_PREFIX = "\n<!--ATP_GENERATED_DOCUMENTS:"
DOCUMENT_MARKER_SUFFIX = ":ATP_GENERATED_DOCUMENTS-->\n"

FORMAT_CONFIG = {
    "pdf": {
        "extension": ".pdf",
        "mime_type": "application/pdf",
        "label": "PDF",
        "icon": "📄",
    },
    "docx": {
        "extension": ".docx",
        "mime_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "label": "Word",
        "icon": "📝",
    },
    "pptx": {
        "extension": ".pptx",
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "label": "PowerPoint",
        "icon": "📊",
    },
    "xlsx": {
        "extension": ".xlsx",
        "mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "label": "Excel",
        "icon": "📈",
    },
    "csv": {
        "extension": ".csv",
        "mime_type": "text/csv",
        "label": "CSV",
        "icon": "📋",
    },
}

_READ_ONLY_PATTERNS = (
    "summarize this pdf", "summarise this pdf", "read this pdf",
    "analyze this pdf", "analyse this pdf", "review this pdf",
    "summarize this word", "read this word", "review this word",
    "analyze this spreadsheet", "review this spreadsheet",
    "summarize this powerpoint", "review this powerpoint",
    "extract from this pdf", "what does this pdf",
)

_FORMAT_PATTERNS = (
    ("pptx", (
        "powerpoint", "power point", "pptx", "slide deck",
        "presentation", "slides",
    )),
    ("xlsx", (
        "excel", "xlsx", "spreadsheet", "workbook",
    )),
    ("docx", (
        "word document", "word file", "docx", "microsoft word",
    )),
    ("csv", (
        "csv", "comma-separated", "comma separated",
    )),
    ("pdf", ("pdf",)),
)

_CREATION_TERMS = (
    "create", "generate", "make", "produce", "prepare", "build", "compile",
    "export", "save", "download", "turn this into", "convert this to",
)


def detect_document_generation_request(prompt_text: Any) -> dict[str, str] | None:
    """
    Detect explicit-format and generic document generation/export requests.

    Generic commands return an empty format so the active workspace can apply
    its saved/default format. Read/review questions remain excluded.
    """
    value = re.sub(r"\s+", " ", str(prompt_text or "")).strip().lower()
    if not value or any(pattern in value for pattern in _READ_ONLY_PATTERNS):
        return None

    requested_format = ""
    for format_name, patterns in _FORMAT_PATTERNS:
        if any(pattern in value for pattern in patterns):
            requested_format = format_name
            break

    # Plain-language follow-ups such as “Convert this to Word” use the
    # standalone product name rather than “Word document” or “DOCX”. Match
    # the complete word only so unrelated text such as “wording” is ignored.
    if not requested_format and re.search(r"\bword\b", value):
        requested_format = "docx"

    creation_terms = (
        "create", "generate", "make", "produce", "prepare", "build", "compile",
        "export", "save", "download", "turn this into", "turn this conversation into",
        "turn this chat into", "convert this to", "convert this into",
        "convert this conversation to", "convert this conversation into",
        "convert the conversation to", "convert the conversation into",
        "convert the current conversation to", "convert the current conversation into",
        "convert the last conversation to", "convert the last conversation into",
        "convert the previous conversation to", "convert the previous conversation into",
        "convert the above conversation to", "convert the above conversation into",
        "convert this chat to", "convert this chat into",
        "convert the current chat to", "convert the current chat into",
        "convert the last chat to", "convert the last chat into",
        "convert the previous chat to", "convert the previous chat into",
        "convert the above chat to", "convert the above chat into",
    )
    artifact_terms = (
        "file", "document", "manual", "guide", "report", "proposal",
        "handbook", "sop", "documentation", "presentation", "slides",
        "spreadsheet", "workbook", "table", "export",
    )
    generic_context_phrases = (
        "convert this to a document",
        "convert this into a document",
        "convert this conversation to a document",
        "convert this conversation into a document",
        "convert the conversation to a document",
        "convert the conversation into a document",
        "convert the current conversation to a document",
        "convert the current conversation into a document",
        "convert the last conversation to a document",
        "convert the last conversation into a document",
        "convert the previous conversation to a document",
        "convert the previous conversation into a document",
        "convert the above conversation to a document",
        "convert the above conversation into a document",
        "convert this chat to a document",
        "convert this chat into a document",
        "convert the current chat to a document",
        "convert the current chat into a document",
        "convert the last chat to a document",
        "convert the last chat into a document",
        "convert the previous chat to a document",
        "convert the previous chat into a document",
        "convert the above chat to a document",
        "convert the above chat into a document",
        "turn this into a document",
        "turn this conversation into a document",
        "turn this chat into a document",
        "export this conversation",
        "export this chat",
        "save this conversation",
        "save this chat",
    )

    creation_requested = any(term in value for term in creation_terms)
    generic_requested = (
        any(phrase in value for phrase in generic_context_phrases)
        or (creation_requested and any(term in value for term in artifact_terms))
    )

    if not requested_format and not generic_requested:
        return None

    if requested_format and not creation_requested:
        if not any(term in value for term in artifact_terms):
            return None

    scope = (
        "conversation"
        if any(
            term in value
            for term in (
                "conversation", "chat",
            )
        )
        else "response"
    )

    if requested_format:
        config = FORMAT_CONFIG[requested_format]
        return {
            "format": requested_format,
            "extension": config["extension"],
            "mime_type": config["mime_type"],
            "label": config["label"],
            "scope": scope,
        }

    return {
        "format": "",
        "extension": "",
        "mime_type": "",
        "label": "Document",
        "scope": "conversation",
    }

def is_document_generation_request(prompt_text: Any) -> bool:
    return detect_document_generation_request(prompt_text) is not None


def _is_document_command_text(value: Any) -> bool:
    text = re.sub(r"\s+", " ", str(value or "")).strip().lower()
    if not text:
        return False
    command_patterns = (
        "convert this to a document",
        "convert this into a document",
        "convert this conversation",
        "convert the conversation",
        "convert the current conversation",
        "convert the last conversation",
        "convert the previous conversation",
        "convert the above conversation",
        "convert this chat",
        "convert the current chat",
        "convert the last chat",
        "convert the previous chat",
        "convert the above chat",
        "turn this into a document",
        "turn this conversation",
        "turn this chat",
        "export this conversation",
        "export this chat",
        "save this conversation",
        "save this chat",
        "create a document",
        "generate a document",
    )
    if any(pattern in text for pattern in command_patterns):
        return True
    return bool(re.search(
        r"\b(?:convert|turn|export|save|create|generate|make|prepare|build|download)\b"
        r".{0,45}\b(?:pdf|docx?|word|pptx?|powerpoint|xlsx?|excel|csv|document|file)\b",
        text,
    ))


def _strip_inline_markdown(value: Any) -> str:
    text = str(value or "")
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"[*_`~]+", "", text)
    return re.sub(r"\s+", " ", text).strip()


def _is_generic_document_heading(value: Any) -> bool:
    text = _strip_inline_markdown(value).casefold().strip(" .:-")
    if not text:
        return True
    generic = {
        "autotecpro ai document", "autotecpro document", "reference document",
        "autotecpro reference document", "product comparison", "comparison document",
        "product comparison document", "document", "report", "guide",
    }
    return text in generic or _is_document_command_text(text)


def _title_from_heading(text: str) -> str:
    headings = re.findall(r"(?m)^\s*#{1,3}\s+(.{4,160}?)\s*$", text)
    for raw in headings:
        heading = _strip_inline_markdown(raw).strip(" .:-")
        if heading and not _is_generic_document_heading(heading):
            return heading[:120]
    return ""


def _model_tokens(text: str) -> list[str]:
    """Return likely product/model identifiers while excluding years and quantities."""
    candidates = re.findall(
        r"(?i)(?<![A-Za-z0-9])(?:[A-Z]{1,6}-)?\d{2,4}(?:[- ]?(?:PRO|S\d+|MK\d+|A\d+))?(?![A-Za-z0-9])",
        text,
    )
    found: list[str] = []
    for raw in candidates:
        token = re.sub(r"\s+", "-", raw.strip())
        if re.fullmatch(r"(?:19|20)\d{2}", token):
            continue
        if token not in found:
            found.append(token)
    return found


def _conversation_subject(answer_text: Any) -> str:
    """Derive a useful subject from headings, conversation questions, and model codes."""
    text = str(answer_text or "").replace("\r\n", "\n").replace("\r", "\n")

    heading = _title_from_heading(text)
    if heading:
        return heading

    speaker_pattern = re.compile(r"(?im)^(User|AutoTecPro AI)\s*:?[ \t]*(.*)$")
    matches = list(speaker_pattern.finditer(text))
    user_candidates: list[str] = []
    for index, match in enumerate(matches):
        if match.group(1).strip().lower() != "user":
            continue
        inline_text = str(match.group(2) or "").strip()
        block_end = matches[index + 1].start() if index + 1 < len(matches) else len(text)
        following_text = text[match.end():block_end].strip()
        candidate = _strip_inline_markdown(" ".join(x for x in (inline_text, following_text) if x))
        if candidate and not _is_document_command_text(candidate):
            user_candidates.append(candidate)

    subject_source = user_candidates[-1] if user_candidates else _strip_inline_markdown(text)
    models = _model_tokens(subject_source)
    if len(models) >= 2:
        return f"AutoTecPro {' vs '.join(models[:3])} Product Comparison"
    if len(models) == 1:
        return f"AutoTecPro Model {models[0]} Reference"

    cleaned = re.sub(
        r"(?i)^\s*(?:please\s+)?(?:can\s+you\s+)?(?:tell\s+me\s+|explain\s+|show\s+me\s+|"
        r"what\s+is\s+|what\s+are\s+|how\s+do\s+i\s+|how\s+to\s+|main\s+details?\s+(?:for|of)\s+)",
        "", subject_source,
    ).strip(" ?.,:-")
    if _is_document_command_text(cleaned):
        cleaned = ""
    words = re.findall(r"[A-Za-z0-9][A-Za-z0-9'&/().+-]*", cleaned)[:14]
    if words:
        title = " ".join(words)
        if not title.casefold().startswith("autotecpro"):
            title = f"AutoTecPro {title}"
        if not re.search(r"(?i)\b(?:document|guide|report|reference|manual|proposal|comparison)\b", title):
            title += " Reference"
        return title[:120].strip()
    return "AutoTecPro Reference Document"


def derive_document_title(prompt_text: Any, answer_text: Any = "") -> str:
    """Return a meaningful title without allowing an export command to become the title."""
    prompt = re.sub(r"\s+", " ", str(prompt_text or "")).strip()
    if prompt and not _is_document_command_text(prompt):
        candidate = _conversation_subject(prompt)
        if candidate and not _is_generic_document_heading(candidate):
            return candidate
    return _conversation_subject(answer_text)


def _safe_stem(prompt_text: Any, answer_text: Any = "") -> str:
    title = derive_document_title(prompt_text, answer_text)
    # Preserve useful model hyphens such as 862-Pro and T732-S3.
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", title)
    cleaned = re.sub(r"_+", "_", cleaned).strip("._-")
    return Path(cleaned[:120] or "AutoTecPro_AI_Document").name


def safe_document_filename(
    prompt_text: Any,
    answer_text: Any = "",
    format_name: str = "pdf",
) -> str:
    format_name = str(format_name or "pdf").lower()
    config = FORMAT_CONFIG.get(format_name, FORMAT_CONFIG["pdf"])
    return _safe_stem(prompt_text, answer_text) + config["extension"]

def _clean_document_text(
    value: Any,
    visible_text_cleaner: Callable[[Any], str] | None = None,
) -> str:
    raw = visible_text_cleaner(value) if callable(visible_text_cleaner) else str(value or "")
    cleaned = str(raw or "").replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"```[A-Za-z0-9_-]*", "", cleaned)
    cleaned = cleaned.replace("```", "")
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip() or "No document content was generated."


def _plain_text(value: Any, cleaner=None) -> str:
    cleaned = _clean_document_text(value, cleaner)
    cleaned = re.sub(r"^\s{0,3}#{1,6}\s*", "", cleaned, flags=re.M)
    cleaned = re.sub(r"\*\*(.*?)\*\*", r"\1", cleaned)
    cleaned = re.sub(r"__(.*?)__", r"\1", cleaned)
    cleaned = re.sub(r"`([^`]+)`", r"\1", cleaned)
    cleaned = cleaned.replace("•", "-")
    return cleaned.strip()


def _pdf_escape_text(value: Any) -> str:
    encoded = str(value or "").encode("cp1252", errors="replace").decode("latin-1")
    return (
        encoded.replace("\\", "\\\\")
        .replace("(", "\\(")
        .replace(")", "\\)")
        .replace("\r", "")
    )


def _wrap_pdf_lines(text: Any, max_chars: int = 92) -> list[str]:
    lines: list[str] = []
    width = max(24, int(max_chars or 92))
    for raw_line in str(text or "").splitlines():
        stripped = raw_line.rstrip()
        if not stripped:
            lines.append("")
            continue
        leading = ""
        content = stripped.strip()
        if content.startswith(("- ", "* ", "• ")):
            leading = "- "
            content = content[2:].strip()
        wrapped = textwrap.wrap(
            content,
            width=max(12, width - len(leading)),
            break_long_words=True,
            break_on_hyphens=True,
            replace_whitespace=False,
            drop_whitespace=True,
        ) or [""]
        for index, line in enumerate(wrapped):
            lines.append((leading if index == 0 else "  ") + line)
    return lines


def build_text_pdf(
    document_text: Any,
    title: str = "AutoTecPro AI Document",
    visible_text_cleaner: Callable[[Any], str] | None = None,
) -> tuple[bytes, int]:
    """Build a valid multi-page PDF without requiring ReportLab."""
    clean_title = re.sub(r"\s+", " ", str(title or "AutoTecPro AI Document")).strip()
    clean_text = _plain_text(document_text, visible_text_cleaner)

    page_width, page_height = 612, 792
    left_margin, top_y, bottom_y = 54, 738, 54
    body_font_size, line_height = 10, 14
    lines_per_page = max(1, int((top_y - bottom_y) / line_height) - 2)
    wrapped_lines = _wrap_pdf_lines(clean_text, max_chars=92)
    pages = [
        wrapped_lines[index:index + lines_per_page]
        for index in range(0, len(wrapped_lines), lines_per_page)
    ] or [["No document content was generated."]]

    objects: dict[int, bytes] = {}
    next_id = 1

    def reserve() -> int:
        nonlocal next_id
        object_id = next_id
        next_id += 1
        return object_id

    catalog_id, pages_id, font_id = reserve(), reserve(), reserve()
    page_ids, content_ids = [], []
    for _ in pages:
        page_ids.append(reserve())
        content_ids.append(reserve())

    objects[font_id] = b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>"

    for page_number, (page_id, content_id, page_lines) in enumerate(
        zip(page_ids, content_ids, pages), start=1
    ):
        commands = [
            "BT", f"/F1 {body_font_size} Tf", f"{left_margin} {top_y} Td",
            f"({_pdf_escape_text(clean_title)}) Tj", f"0 -{line_height * 2} Td",
        ]
        for line in page_lines:
            commands.extend([
                f"({_pdf_escape_text(line)}) Tj",
                f"0 -{line_height} Td",
            ])
        footer = f"Page {page_number} of {len(pages)}"
        commands.extend([
            "ET", "BT", "/F1 8 Tf", f"{page_width - 110} 30 Td",
            f"({_pdf_escape_text(footer)}) Tj", "ET",
        ])
        stream = "\n".join(commands).encode("latin-1", errors="replace")
        objects[content_id] = (
            f"<< /Length {len(stream)} >>\nstream\n".encode("ascii")
            + stream + b"\nendstream"
        )
        objects[page_id] = (
            f"<< /Type /Page /Parent {pages_id} 0 R "
            f"/MediaBox [0 0 {page_width} {page_height}] "
            f"/Resources << /Font << /F1 {font_id} 0 R >> >> "
            f"/Contents {content_id} 0 R >>"
        ).encode("ascii")

    kids = " ".join(f"{page_id} 0 R" for page_id in page_ids)
    objects[pages_id] = (
        f"<< /Type /Pages /Kids [{kids}] /Count {len(page_ids)} >>"
    ).encode("ascii")
    objects[catalog_id] = (
        f"<< /Type /Catalog /Pages {pages_id} 0 R >>"
    ).encode("ascii")

    pdf = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0] * next_id
    for object_id in range(1, next_id):
        offsets[object_id] = len(pdf)
        pdf.extend(f"{object_id} 0 obj\n".encode("ascii"))
        pdf.extend(objects[object_id])
        pdf.extend(b"\nendobj\n")
    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {next_id}\n".encode("ascii"))
    pdf.extend(b"0000000000 65535 f \n")
    for object_id in range(1, next_id):
        pdf.extend(f"{offsets[object_id]:010d} 00000 n \n".encode("ascii"))
    pdf.extend(
        (
            f"trailer\n<< /Size {next_id} /Root {catalog_id} 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF"
        ).encode("ascii")
    )
    return bytes(pdf), len(pages)


def _strip_docx_inline_markdown(value: Any) -> str:
    """Remove lightweight inline Markdown while preserving readable content."""
    text = str(value or "")
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"__(.*?)__", r"\1", text)
    text = re.sub(r"(?<!\*)\*([^*]+)\*(?!\*)", r"\1", text)
    return text.strip()


def _split_docx_markdown_table_row(value: Any) -> list[str]:
    """Split a Markdown table row, tolerating optional outer pipes."""
    line = str(value or "").strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|") and not line.endswith(r"\|"):
        line = line[:-1]
    placeholder = "\x00ATP_ESCAPED_PIPE\x00"
    line = line.replace(r"\|", placeholder)
    cells = [
        _strip_docx_inline_markdown(cell.replace(placeholder, "|").strip())
        for cell in line.split("|")
    ]
    return cells


def _is_docx_markdown_table_separator(value: Any) -> bool:
    cells = _split_docx_markdown_table_row(value)
    return bool(cells) and all(
        re.fullmatch(r":?-{3,}:?", str(cell or "").replace(" ", ""))
        for cell in cells
    )


def _looks_like_markdown_table_row(value: Any) -> bool:
    line = str(value or "").strip()
    return "|" in line and len(_split_docx_markdown_table_row(line)) >= 2


def _iter_markdown_blocks(text: str):
    """Yield grouped Markdown blocks, including native-table candidates."""
    lines = str(text or "").splitlines()
    index = 0
    while index < len(lines):
        raw = lines[index]
        stripped = raw.strip()

        # A valid Markdown table starts with a header row followed by a divider.
        if (
            _looks_like_markdown_table_row(stripped)
            and index + 1 < len(lines)
            and _is_docx_markdown_table_separator(lines[index + 1])
        ):
            rows = [_split_docx_markdown_table_row(stripped)]
            index += 2  # Skip the Markdown divider; it is formatting, not content.
            while index < len(lines) and _looks_like_markdown_table_row(lines[index]):
                if not _is_docx_markdown_table_separator(lines[index]):
                    rows.append(_split_docx_markdown_table_row(lines[index]))
                index += 1
            width = max((len(row) for row in rows), default=1)
            normalized = [row + [""] * (width - len(row)) for row in rows]
            yield ("table", normalized)
            continue

        if not stripped:
            yield ("blank", "")
        elif re.match(r"^#{1,6}\s+", stripped):
            level = len(stripped) - len(stripped.lstrip("#"))
            yield (f"heading{min(level, 3)}", stripped[level:].strip())
        elif stripped.startswith(">"):
            yield ("quote", stripped.lstrip(">").strip())
        elif stripped.startswith(("- ", "* ", "• ")):
            yield ("bullet", stripped[2:].strip())
        elif re.match(r"^\d+[.)]\s+", stripped):
            yield ("number", re.sub(r"^\d+[.)]\s+", "", stripped))
        else:
            yield ("paragraph", stripped)
        index += 1


def _add_docx_inline_text(paragraph, value: Any, *, bold: bool = False) -> None:
    """Add basic bold/code Markdown as Word runs instead of exposing markers."""
    text = str(value or "")
    token_pattern = re.compile(r"(\*\*.*?\*\*|__.*?__|`[^`]+`)")
    cursor = 0
    for match in token_pattern.finditer(text):
        if match.start() > cursor:
            run = paragraph.add_run(text[cursor:match.start()])
            run.bold = bold
        token = match.group(0)
        if token.startswith("`"):
            run = paragraph.add_run(token[1:-1])
            run.bold = bold
            run.font.name = "Consolas"
        else:
            run = paragraph.add_run(token[2:-2])
            run.bold = True
        cursor = match.end()
    if cursor < len(text):
        run = paragraph.add_run(text[cursor:])
        run.bold = bold


def _set_docx_cell_text(cell, value: Any, *, bold: bool = False) -> None:
    cell.text = ""
    paragraph = cell.paragraphs[0]
    paragraph.paragraph_format.space_after = 0
    _add_docx_inline_text(paragraph, value, bold=bold)


def _shade_docx_element(element, fill: str) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    xml_element = getattr(element, "_tc", None)
    if xml_element is None:
        xml_element = getattr(element, "_p", None)
    if xml_element is None:
        xml_element = element
    properties = (
        xml_element.get_or_add_tcPr()
        if hasattr(xml_element, "get_or_add_tcPr")
        else xml_element.get_or_add_pPr()
    )
    shade = properties.find(qn("w:shd"))
    if shade is None:
        shade = OxmlElement("w:shd")
        properties.append(shade)
    shade.set(qn("w:fill"), fill)


def _configure_docx_styles(document) -> None:
    from docx.enum.text import WD_LINE_SPACING
    from docx.shared import Pt, RGBColor

    styles = document.styles
    normal = styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(10.5)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing_rule = WD_LINE_SPACING.SINGLE

    heading_settings = {
        "Title": (20, 14),
        "Heading 1": (15, 12),
        "Heading 2": (12.5, 10),
        "Heading 3": (11, 8),
    }
    for style_name, (size, before) in heading_settings.items():
        if style_name not in styles:
            continue
        style = styles[style_name]
        style.font.name = "Arial"
        style.font.size = Pt(size)
        style.font.bold = True
        style.font.color.rgb = RGBColor(128, 24, 24)
        style.paragraph_format.space_before = Pt(before)
        style.paragraph_format.space_after = Pt(5)
        style.paragraph_format.keep_with_next = True

    for style_name in ("List Bullet", "List Number"):
        if style_name in styles:
            styles[style_name].font.name = "Arial"
            styles[style_name].font.size = Pt(10.5)
            styles[style_name].paragraph_format.space_after = Pt(2)


def _add_docx_markdown_table(document, rows: list[list[str]]) -> None:
    from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
    from docx.shared import Pt

    if not rows:
        return
    column_count = max(len(row) for row in rows)
    table = document.add_table(rows=len(rows), cols=column_count)
    table.style = "Table Grid"
    table.autofit = True

    for row_index, row in enumerate(rows):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            value = row[column_index] if column_index < len(row) else ""
            _set_docx_cell_text(cell, value, bold=(row_index == 0))
            for paragraph in cell.paragraphs:
                paragraph.paragraph_format.space_before = Pt(0)
                paragraph.paragraph_format.space_after = Pt(2)
            if row_index == 0:
                _shade_docx_element(cell, "D9E2F3")

    document.add_paragraph("").paragraph_format.space_after = Pt(2)


def build_docx(
    document_text: Any,
    title: str,
    visible_text_cleaner=None,
) -> tuple[bytes, int]:
    """Build a polished DOCX with native headings, lists, quotes, and tables."""
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt

    text = _clean_document_text(document_text, visible_text_cleaner)
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.7)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.8)
    section.right_margin = Inches(0.8)

    _configure_docx_styles(document)
    title_paragraph = document.add_heading(_strip_docx_inline_markdown(title), 0)
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

    previous_kind = ""
    first_content_block = True
    normalized_title = _strip_docx_inline_markdown(title).casefold().strip(" .:-")
    for kind, value in _iter_markdown_blocks(text):
        if (
            first_content_block
            and kind.startswith("heading")
            and _strip_docx_inline_markdown(value).casefold().strip(" .:-") == normalized_title
        ):
            # The document title already appears above; avoid printing it twice.
            first_content_block = False
            previous_kind = kind
            continue
        if kind != "blank":
            first_content_block = False
        if kind == "blank":
            # Avoid large empty gaps while retaining separation between sections.
            if previous_kind not in {"", "blank", "heading1", "heading2", "heading3", "table"}:
                spacer = document.add_paragraph("")
                spacer.paragraph_format.space_after = Pt(1)
            previous_kind = "blank"
            continue

        if kind == "table":
            _add_docx_markdown_table(document, value)
        elif kind.startswith("heading"):
            paragraph = document.add_heading(
                _strip_docx_inline_markdown(value),
                level=int(kind[-1]),
            )
        elif kind == "bullet":
            paragraph = document.add_paragraph(style="List Bullet")
            _add_docx_inline_text(paragraph, value)
        elif kind == "number":
            paragraph = document.add_paragraph(style="List Number")
            _add_docx_inline_text(paragraph, value)
        elif kind == "quote":
            style_name = "Intense Quote" if "Intense Quote" in document.styles else "Quote"
            paragraph = document.add_paragraph(style=style_name)
            _add_docx_inline_text(paragraph, value)
            paragraph.paragraph_format.space_before = Pt(3)
            paragraph.paragraph_format.space_after = Pt(3)
            _shade_docx_element(paragraph._p, "F2F2F2")
        else:
            clean_value = str(value or "").strip()
            # Conversation speaker labels should read as labels, not body copy.
            if clean_value.casefold() in {"user", "autotecpro ai", "assistant"}:
                paragraph = document.add_paragraph()
                paragraph.paragraph_format.space_before = Pt(8)
                paragraph.paragraph_format.space_after = Pt(2)
                _add_docx_inline_text(paragraph, clean_value, bold=True)
            else:
                paragraph = document.add_paragraph()
                _add_docx_inline_text(paragraph, clean_value)

        previous_kind = kind

    buffer = io.BytesIO()
    document.save(buffer)
    element_count = len(document.paragraphs) + len(document.tables)
    return buffer.getvalue(), max(1, element_count)

def _split_slide_sections(text: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    current_lines: list[str] = []
    for kind, value in _iter_markdown_blocks(text):
        if kind.startswith("heading"):
            if current_title or current_lines:
                sections.append((current_title or "Overview", current_lines))
            current_title, current_lines = value, []
        elif kind != "blank":
            current_lines.append(value)
    if current_title or current_lines:
        sections.append((current_title or "Overview", current_lines))
    return sections or [("Overview", [text])]


def build_pptx(
    document_text: Any,
    title: str,
    visible_text_cleaner=None,
) -> tuple[bytes, int]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    text = _clean_document_text(document_text, visible_text_cleaner)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Generated by AutoTecPro AI"

    for section_title, section_lines in _split_slide_sections(text):
        chunks = [
            section_lines[index:index + 7]
            for index in range(0, len(section_lines), 7)
        ] or [[]]
        for chunk_index, chunk in enumerate(chunks):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = (
                section_title if chunk_index == 0
                else f"{section_title} (continued)"
            )
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for item_index, line in enumerate(chunk):
                paragraph = frame.paragraphs[0] if item_index == 0 else frame.add_paragraph()
                paragraph.text = re.sub(r"\*\*(.*?)\*\*", r"\1", line)
                paragraph.level = 0
                paragraph.font.size = Pt(22)

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue(), len(presentation.slides)


def _extract_table_rows(text: str) -> list[list[str]]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    markdown_rows: list[list[str]] = []
    for line in lines:
        if "|" in line:
            cells = [cell.strip() for cell in line.strip("|").split("|")]
            if cells and not all(re.fullmatch(r":?-{3,}:?", c or "") for c in cells):
                markdown_rows.append(cells)
    if len(markdown_rows) >= 2:
        width = max(len(row) for row in markdown_rows)
        return [row + [""] * (width - len(row)) for row in markdown_rows]

    rows = []
    for line in lines:
        line = re.sub(r"^[-*•]\s+", "", line)
        if ":" in line and len(line.split(":", 1)[0]) <= 50:
            left, right = line.split(":", 1)
            rows.append([left.strip(), right.strip()])
        elif "," in line:
            rows.append([cell.strip() for cell in next(csv.reader([line]))])
        else:
            rows.append([line])
    width = max((len(row) for row in rows), default=1)
    return [row + [""] * (width - len(row)) for row in rows]


def build_xlsx(
    document_text: Any,
    title: str,
    visible_text_cleaner=None,
) -> tuple[bytes, int]:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    text = _clean_document_text(document_text, visible_text_cleaner)
    rows = _extract_table_rows(text)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Document"
    sheet["A1"] = title
    sheet["A1"].font = Font(bold=True, size=16)
    sheet.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(1, len(rows[0]) if rows else 1))

    for r_index, row in enumerate(rows, start=3):
        for c_index, value in enumerate(row, start=1):
            cell = sheet.cell(row=r_index, column=c_index, value=value)
            cell.alignment = Alignment(vertical="top", wrap_text=True)
            if r_index == 3 and len(rows) > 1:
                cell.font = Font(bold=True)
                cell.fill = PatternFill("solid", fgColor="D9EAF7")

    for c_index in range(1, max((len(r) for r in rows), default=1) + 1):
        max_len = max(
            [len(str(sheet.cell(r, c_index).value or "")) for r in range(1, sheet.max_row + 1)]
            or [10]
        )
        sheet.column_dimensions[get_column_letter(c_index)].width = min(max(max_len + 2, 12), 40)
    sheet.freeze_panes = "A3"
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue(), max(1, len(workbook.sheetnames))


def build_csv(
    document_text: Any,
    title: str,
    visible_text_cleaner=None,
) -> tuple[bytes, int]:
    text = _clean_document_text(document_text, visible_text_cleaner)
    rows = _extract_table_rows(text)
    buffer = io.StringIO(newline="")
    writer = csv.writer(buffer)
    writer.writerows(rows)
    return buffer.getvalue().encode("utf-8-sig"), max(1, len(rows))


def serialize_documents_marker(documents: Iterable[dict[str, Any]] | None) -> str:
    records = list(documents or [])
    if not records:
        return ""
    try:
        return (
            "\n\n" + DOCUMENT_MARKER_PREFIX
            + json.dumps(records, ensure_ascii=False)
            + DOCUMENT_MARKER_SUFFIX
        )
    except Exception:
        return ""


def extract_documents_from_message_content(
    content: Any,
) -> tuple[str, list[dict[str, Any]]]:
    value = str(content or "")
    pattern = re.escape(DOCUMENT_MARKER_PREFIX) + r"(.*?)" + re.escape(DOCUMENT_MARKER_SUFFIX)
    match = re.search(pattern, value, flags=re.DOTALL)
    if not match:
        # Backward compatibility with PDF-only history markers.
        legacy_pattern = (
            re.escape("\n<!--ATP_PDF_DOCUMENTS:")
            + r"(.*?)"
            + re.escape(":ATP_PDF_DOCUMENTS-->\n")
        )
        match = re.search(legacy_pattern, value, flags=re.DOTALL)
        if not match:
            return value, []

    visible_text = (value[:match.start()] + value[match.end():]).strip()
    try:
        documents = json.loads(match.group(1))
    except Exception:
        documents = []

    clean_documents: list[dict[str, Any]] = []
    for document in documents if isinstance(documents, list) else []:
        if not isinstance(document, dict):
            continue
        data_url = str(document.get("data_url") or "")
        if not data_url.startswith("data:") or ";base64," not in data_url:
            continue
        format_name = str(document.get("format") or "").lower()
        if not format_name:
            mime = str(document.get("mime_type") or data_url[5:].split(";", 1)[0])
            format_name = next(
                (key for key, config in FORMAT_CONFIG.items() if config["mime_type"] == mime),
                "pdf",
            )
        if format_name not in FORMAT_CONFIG:
            continue
        config = FORMAT_CONFIG[format_name]
        try:
            size_bytes = int(document.get("size_bytes") or 0)
        except (TypeError, ValueError):
            size_bytes = 0
        try:
            unit_count = int(document.get("unit_count") or document.get("page_count") or 0)
        except (TypeError, ValueError):
            unit_count = 0
        clean_documents.append({
            "name": Path(str(document.get("name") or ("AutoTecPro_AI_Document" + config["extension"]))).name,
            "format": format_name,
            "format_label": document.get("format_label") or config["label"],
            "icon": document.get("icon") or config["icon"],
            "mime_type": config["mime_type"],
            "data_url": data_url,
            "unit_count": unit_count,
            "unit_label": document.get("unit_label") or ("Pages" if format_name == "pdf" else ""),
            "page_count": unit_count if format_name == "pdf" else 0,
            "size_bytes": size_bytes,
            "generated": bool(document.get("generated", True)),
            "created_at": document.get("created_at"),
        })
    return visible_text, clean_documents

# ============================================================
# Contextual document export and optional branding overrides
# ============================================================

def _document_options(options: Any = None) -> dict[str, Any]:
    raw = dict(options or {}) if isinstance(options, dict) else {}
    style = str(raw.get("style") or "Professional").strip().title()
    if style not in {"Professional", "Minimal", "Presentation"}:
        style = "Professional"
    logo_path = Path(str(raw.get("logo_path") or "")).expanduser()
    return {
        "include_logo": bool(raw.get("include_logo")),
        "include_company_info": bool(raw.get("include_company_info")),
        "watermark": str(raw.get("watermark") or "").strip()[:60],
        "style": style,
        "logo_path": logo_path if logo_path.is_file() else None,
        "company_name": str(raw.get("company_name") or "AutoTecPro Inc.").strip(),
        "company_info": str(raw.get("company_info") or "Generated by AutoTecPro AI").strip(),
    }



def _prepared_logo_stream(logo_path: Any):
    """
    Return a tightly cropped PNG stream for document embedding.

    The AutoTecPro source logo uses a large opaque white canvas. Prefer the
    visible non-white artwork bounds; only use alpha bounds when transparency
    actually removes part of the image.
    """
    if not logo_path:
        return None

    path = Path(str(logo_path))
    if not path.exists() or not path.is_file():
        return None

    try:
        from PIL import Image, ImageChops

        image = Image.open(path).convert("RGBA")
        original_size = image.size
        alpha = image.getchannel("A")
        alpha_box = alpha.getbbox()

        rgb_on_white = Image.new("RGB", image.size, "white")
        rgb_on_white.paste(image.convert("RGB"), mask=alpha)
        white_background = Image.new("RGB", rgb_on_white.size, "white")
        difference = ImageChops.difference(
            rgb_on_white,
            white_background,
        ).convert("L")
        visible_mask = difference.point(
            lambda pixel: 255 if pixel > 12 else 0
        )
        color_box = visible_mask.getbbox()

        crop_box = None
        if color_box:
            crop_box = color_box
        elif alpha_box and alpha_box != (0, 0, image.width, image.height):
            crop_box = alpha_box

        if crop_box:
            left, top, right, bottom = crop_box
            padding = max(6, int(min(original_size) * 0.012))
            crop_box = (
                max(0, left - padding),
                max(0, top - padding),
                min(image.width, right + padding),
                min(image.height, bottom + padding),
            )
            image = image.crop(crop_box)

        # Flatten to RGB so PDF/Office renderers handle the image consistently.
        flattened = Image.new("RGB", image.size, "white")
        flattened.paste(image.convert("RGB"), mask=image.getchannel("A"))

        stream = io.BytesIO()
        flattened.save(stream, format="PNG", optimize=True)
        stream.seek(0)
        return stream
    except Exception:
        return None

def _logo_dimensions(logo_stream, max_width, max_height):
    """Preserve the logo aspect ratio within the requested bounds."""
    try:
        from PIL import Image
        position = logo_stream.tell()
        image = Image.open(logo_stream)
        width, height = image.size
        logo_stream.seek(position)
        if width <= 0 or height <= 0:
            return max_width, max_height
        scale = min(max_width / width, max_height / height)
        return width * scale, height * scale
    except Exception:
        return max_width, max_height

def _style_profile(style: str) -> dict[str, Any]:
    """Return materially different layout rules for each document style."""
    clean = str(style or "Professional").strip().title()
    profiles = {
        "Professional": {
            "page_margins": (0.72, 0.72, 0.68, 0.68),
            "title_size": 23,
            "title_leading": 28,
            "body_size": 10.5,
            "body_leading": 15,
            "heading_size": 15,
            "heading_leading": 19,
            "cover": False,
            "title_alignment": "center",
            "accent": "#B91C1C",
            "body_color": "#111827",
            "page_numbers": True,
            "logo_width": 2.15,
            "logo_height": 0.78,
            "content_spacing": 7,
        },
        "Minimal": {
            "page_margins": (0.58, 0.58, 0.52, 0.52),
            "title_size": 17,
            "title_leading": 21,
            "body_size": 9.5,
            "body_leading": 12.5,
            "heading_size": 12,
            "heading_leading": 15,
            "cover": False,
            "title_alignment": "left",
            "accent": "#1F2937",
            "body_color": "#111827",
            "page_numbers": False,
            "logo_width": 1.35,
            "logo_height": 0.48,
            "content_spacing": 3,
        },
        "Presentation": {
            "page_margins": (0.92, 0.92, 0.82, 0.82),
            "title_size": 30,
            "title_leading": 36,
            "body_size": 13,
            "body_leading": 20,
            "heading_size": 20,
            "heading_leading": 25,
            "cover": False,
            "title_alignment": "center",
            "accent": "#B91C1C",
            "body_color": "#172033",
            "page_numbers": True,
            "logo_width": 2.55,
            "logo_height": 0.92,
            "content_spacing": 11,
        },
    }
    return profiles.get(clean, profiles["Professional"])


def _draw_pdf_watermark(canvas, page_width, page_height, text, colors):
    """Draw a large diagonal watermark that is clearly visible but readable."""
    if not text:
        return
    canvas.saveState()
    try:
        canvas.setFillAlpha(0.13)
    except Exception:
        pass
    canvas.setFillColor(colors.HexColor("#6B7280"))
    canvas.setFont("Helvetica-Bold", 48)
    canvas.translate(page_width / 2, page_height / 2)
    canvas.rotate(35)
    canvas.drawCentredString(0, 0, str(text).upper())
    canvas.restoreState()


def _split_markdown_table_row(line: str) -> list[str]:
    value = line.strip().strip("|")
    # Support escaped pipes inside cells.
    parts = re.split(r"(?<!\\)\|", value)
    return [_strip_inline_markdown(part.replace(r"\|", "|")).strip() for part in parts]


def _is_markdown_table_separator(line: str) -> bool:
    cells = _split_markdown_table_row(line)
    return bool(cells) and all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells)


def _iter_document_blocks(text: str):
    """Yield paragraphs/headings/lists and complete Markdown tables."""
    lines = text.splitlines()
    i = 0
    while i < len(lines):
        raw = lines[i]
        stripped = raw.strip()
        if (
            "|" in stripped
            and i + 1 < len(lines)
            and _is_markdown_table_separator(lines[i + 1].strip())
        ):
            rows = [_split_markdown_table_row(stripped)]
            i += 2
            while i < len(lines):
                candidate = lines[i].strip()
                if not candidate or "|" not in candidate:
                    break
                rows.append(_split_markdown_table_row(candidate))
                i += 1
            width = max((len(row) for row in rows), default=1)
            yield "table", [row + [""] * (width - len(row)) for row in rows]
            continue
        if not stripped:
            yield "blank", ""
        elif re.match(r"^#{1,6}\s+", stripped):
            level = len(stripped) - len(stripped.lstrip("#"))
            yield f"heading{min(level, 3)}", stripped[level:].strip()
        elif stripped.startswith(("- ", "* ", "• ")):
            yield "bullet", stripped[2:].strip()
        elif re.match(r"^\d+[.)]\s+", stripped):
            yield "number", stripped
        else:
            yield "paragraph", stripped
        i += 1


def _pdf_table_column_widths(rows, available_width):
    count = max((len(row) for row in rows), default=1)
    if count <= 1:
        return [available_width]
    lengths = []
    for col in range(count):
        values = [str(row[col] if col < len(row) else "") for row in rows]
        lengths.append(max(8, min(45, max((len(v) for v in values), default=8))))
    total = sum(lengths) or count
    widths = [available_width * length / total for length in lengths]
    minimum = min(0.85 * 72, available_width / count)
    widths = [max(minimum, width) for width in widths]
    scale = available_width / sum(widths)
    return [width * scale for width in widths]


def _build_pdf_with_options(document_text, title, cleaner=None, options=None):
    opts = _document_options(options)
    profile = _style_profile(opts["style"])
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from reportlab.lib import colors
        from reportlab.lib.units import inch
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, PageBreak,
            Image as RLImage, KeepTogether, HRFlowable, LongTable, TableStyle,
        )
        import html as _html
    except Exception:
        return build_text_pdf(document_text, title, cleaner)

    text = _clean_document_text(document_text, cleaner)
    buffer = io.BytesIO()
    right, left, top, bottom = profile["page_margins"]
    pdf_doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=right * inch,
        leftMargin=left * inch,
        topMargin=top * inch,
        bottomMargin=bottom * inch,
        title=title,
        author=opts["company_name"],
        subject=f"{opts['style']} AutoTecPro document",
    )

    styles = getSampleStyleSheet()
    accent = colors.HexColor(profile["accent"])
    body_color = colors.HexColor(profile["body_color"])
    alignment = TA_CENTER if profile["title_alignment"] == "center" else TA_LEFT

    styles.add(ParagraphStyle(
        name="ATPTitleV2",
        parent=styles["Title"],
        textColor=accent,
        fontName="Helvetica-Bold",
        fontSize=profile["title_size"],
        leading=profile["title_leading"],
        alignment=alignment,
        spaceAfter=18 if profile["cover"] else 10,
    ))
    styles.add(ParagraphStyle(
        name="ATPSubtitleV2",
        parent=styles["BodyText"],
        textColor=colors.HexColor("#6B7280"),
        fontName="Helvetica",
        fontSize=9.5 if opts["style"] != "Presentation" else 11,
        leading=13,
        alignment=alignment,
        spaceAfter=15,
    ))
    styles.add(ParagraphStyle(
        name="ATPBodyV2",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=profile["body_size"],
        leading=profile["body_leading"],
        textColor=body_color,
        spaceAfter=profile["content_spacing"],
    ))
    styles.add(ParagraphStyle(
        name="ATPH1V2",
        parent=styles["Heading1"],
        textColor=accent,
        fontName="Helvetica-Bold",
        fontSize=profile["heading_size"],
        leading=profile["heading_leading"],
        spaceBefore=14 if opts["style"] != "Minimal" else 8,
        spaceAfter=8 if opts["style"] != "Minimal" else 4,
    ))
    styles.add(ParagraphStyle(
        name="ATPBulletV2",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=profile["body_size"],
        leading=profile["body_leading"],
        textColor=body_color,
        leftIndent=13 if opts["style"] != "Presentation" else 18,
        firstLineIndent=-8,
        bulletIndent=2,
        spaceAfter=profile["content_spacing"],
    ))

    story = []
    logo_added = False
    if opts["include_logo"]:
        logo_stream = _prepared_logo_stream(opts["logo_path"])
        if logo_stream is None:
            raise RuntimeError(
                "AutoTecPro logo was selected, but App/logo.png could not be loaded."
            )
        logo_width, logo_height = _logo_dimensions(
            logo_stream,
            max_width=profile["logo_width"] * inch,
            max_height=profile["logo_height"] * inch,
        )
        logo = RLImage(logo_stream, width=logo_width, height=logo_height)
        logo.hAlign = "CENTER" if profile["title_alignment"] == "center" else "LEFT"
        story.extend([logo, Spacer(1, 10 if opts["style"] != "Minimal" else 5)])
        logo_added = True

    story.append(Paragraph(_html.escape(title), styles["ATPTitleV2"]))

    if opts["include_company_info"]:
        company_lines = [
            _html.escape(opts["company_name"]),
            _html.escape(opts["company_info"]),
        ]
        story.append(Paragraph("<br/>".join(company_lines), styles["ATPSubtitleV2"]))

    # Keep branding compactly on the first content page. No cover page.
    story.append(HRFlowable(
        width="100%" if opts["style"] == "Minimal" else "72%",
        thickness=0.6 if opts["style"] == "Minimal" else 1.1,
        color=(
            colors.HexColor("#D1D5DB")
            if opts["style"] == "Minimal"
            else accent
        ),
        hAlign=(
            "LEFT"
            if profile["title_alignment"] == "left"
            else "CENTER"
        ),
        spaceBefore=1,
        spaceAfter=8 if opts["style"] != "Presentation" else 12,
    ))

    for kind, value in _iter_document_blocks(text):
        if kind == "table":
            table_rows = value
            cell_style = ParagraphStyle(
                name=f"ATPTableCell{len(story)}",
                parent=styles["ATPBodyV2"],
                fontSize=max(7.4, profile["body_size"] - 1.8),
                leading=max(9.2, profile["body_leading"] - 3.2),
                spaceAfter=0,
                textColor=body_color,
            )
            header_style = ParagraphStyle(
                name=f"ATPTableHeader{len(story)}",
                parent=cell_style,
                fontName="Helvetica-Bold",
                textColor=colors.white,
            )
            formatted = []
            for row_index, row in enumerate(table_rows):
                formatted.append([
                    Paragraph(_html.escape(str(cell)), header_style if row_index == 0 else cell_style)
                    for cell in row
                ])
            widths = _pdf_table_column_widths(table_rows, pdf_doc.width)
            table = LongTable(
                formatted,
                colWidths=widths,
                repeatRows=1,
                hAlign="LEFT",
                splitByRow=1,
            )
            commands = [
                ("BACKGROUND", (0, 0), (-1, 0), accent),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("GRID", (0, 0), (-1, -1), 0.45, colors.HexColor("#CBD5E1")),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]
            for row_index in range(1, len(formatted)):
                if row_index % 2 == 0:
                    commands.append(("BACKGROUND", (0, row_index), (-1, row_index), colors.HexColor("#F8FAFC")))
            table.setStyle(TableStyle(commands))
            story.extend([Spacer(1, 5), table, Spacer(1, 9)])
            continue

        clean = re.sub(r"\*\*(.*?)\*\*", r"\1", value)
        clean = re.sub(r"`([^`]+)`", r"\1", clean)
        safe = _html.escape(clean)
        if kind == "blank":
            story.append(Spacer(1, 3 if opts["style"] == "Minimal" else 7))
        elif kind.startswith("heading"):
            # Avoid repeating the generated title as the first body heading.
            if _strip_inline_markdown(clean).casefold() != _strip_inline_markdown(title).casefold():
                story.append(Paragraph(safe, styles["ATPH1V2"]))
        elif kind == "bullet":
            story.append(Paragraph("• " + safe, styles["ATPBulletV2"]))
        elif kind == "number":
            story.append(Paragraph(safe, styles["ATPBulletV2"]))
        else:
            story.append(Paragraph(safe, styles["ATPBodyV2"]))

    def draw_page(canvas, doc_obj):
        canvas.saveState()
        _draw_pdf_watermark(
            canvas,
            letter[0],
            letter[1],
            opts["watermark"],
            colors,
        )

        if opts["include_company_info"]:
            canvas.setFillColor(colors.HexColor("#6B7280"))
            canvas.setFont("Helvetica", 7.8)
            canvas.drawString(
                left * inch,
                0.30 * inch,
                f"{opts['company_name']} • {opts['company_info']}",
            )

        if profile["page_numbers"]:
            canvas.setFillColor(colors.HexColor("#6B7280"))
            canvas.setFont("Helvetica", 8)
            canvas.drawRightString(
                letter[0] - right * inch,
                0.30 * inch,
                f"Page {doc_obj.page}",
            )
        canvas.restoreState()

    pdf_doc.build(
        story,
        onFirstPage=draw_page,
        onLaterPages=draw_page,
    )
    data = buffer.getvalue()
    page_count = len(re.findall(rb"/Type\s*/Page\b", data))
    return data, max(1, page_count)


def _set_docx_cell_watermark(section, text):
    """Create a clear repeating confidential header for Word documents."""
    if not text:
        return
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    header = section.header
    paragraph = header.paragraphs[0]
    paragraph.text = str(text).upper()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.runs[0]
    run.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = RGBColor(180, 180, 180)


def _build_docx_with_options(document_text, title, cleaner=None, options=None):
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    opts = _document_options(options)
    profile = _style_profile(opts["style"])
    text = _clean_document_text(document_text, cleaner)
    document = Document()
    section = document.sections[0]

    right, left, top, bottom = profile["page_margins"]
    section.top_margin = Inches(top)
    section.bottom_margin = Inches(bottom)
    section.left_margin = Inches(left)
    section.right_margin = Inches(right)

    normal = document.styles["Normal"]
    normal.font.name = "Arial"
    normal.font.size = Pt(profile["body_size"])

    title_style = document.styles["Title"]
    title_style.font.name = "Arial"
    title_style.font.size = Pt(profile["title_size"])
    title_style.font.bold = True
    title_style.font.color.rgb = (
        RGBColor(185, 28, 28)
        if opts["style"] != "Minimal"
        else RGBColor(31, 41, 55)
    )

    for level in (1, 2, 3):
        heading_style = document.styles[f"Heading {level}"]
        heading_style.font.name = "Arial"
        heading_style.font.bold = True
        heading_style.font.size = Pt(
            max(11, profile["heading_size"] - ((level - 1) * 2))
        )
        heading_style.font.color.rgb = (
            RGBColor(185, 28, 28)
            if opts["style"] != "Minimal"
            else RGBColor(31, 41, 55)
        )

    logo_added = False
    if opts["include_logo"]:
        logo_stream = _prepared_logo_stream(opts["logo_path"])
        if logo_stream is None:
            raise RuntimeError(
                "AutoTecPro logo was selected, but App/logo.png could not be loaded."
            )
        p = document.add_paragraph()
        p.alignment = (
            WD_ALIGN_PARAGRAPH.CENTER
            if profile["title_alignment"] == "center"
            else WD_ALIGN_PARAGRAPH.LEFT
        )
        p.add_run().add_picture(
            logo_stream,
            width=Inches(profile["logo_width"]),
        )
        logo_added = True

    heading = document.add_heading(title, 0)
    heading.alignment = (
        WD_ALIGN_PARAGRAPH.CENTER
        if profile["title_alignment"] == "center"
        else WD_ALIGN_PARAGRAPH.LEFT
    )

    if opts["include_company_info"]:
        p = document.add_paragraph()
        p.alignment = heading.alignment
        run = p.add_run(
            f"{opts['company_name']}\n{opts['company_info']}"
        )
        run.font.size = Pt(9.5 if opts["style"] != "Presentation" else 11)
        run.font.color.rgb = RGBColor(107, 114, 128)

    previous_kind = ""
    first_content_block = True
    normalized_title = _strip_docx_inline_markdown(title).casefold().strip(" .:-")

    for kind, value in _iter_markdown_blocks(text):
        if (
            first_content_block
            and kind.startswith("heading")
            and _strip_docx_inline_markdown(value).casefold().strip(" .:-")
            == normalized_title
        ):
            first_content_block = False
            previous_kind = kind
            continue

        if kind != "blank":
            first_content_block = False

        if kind == "table":
            _add_docx_markdown_table(document, value)
            previous_kind = kind
            continue

        if kind == "blank":
            if previous_kind not in {
                "", "blank", "heading1", "heading2", "heading3", "table"
            }:
                paragraph = document.add_paragraph("")
                paragraph.paragraph_format.space_after = Pt(1)
            previous_kind = "blank"
            continue

        clean_value = _strip_docx_inline_markdown(value)

        if kind.startswith("heading"):
            paragraph = document.add_heading(
                clean_value,
                level=int(kind[-1]),
            )
        elif kind == "bullet":
            paragraph = document.add_paragraph(style="List Bullet")
            _add_docx_inline_text(paragraph, value)
        elif kind == "number":
            paragraph = document.add_paragraph(style="List Number")
            _add_docx_inline_text(paragraph, value)
        elif kind == "quote":
            style_name = (
                "Intense Quote"
                if "Intense Quote" in document.styles
                else "Quote"
            )
            paragraph = document.add_paragraph(style=style_name)
            _add_docx_inline_text(paragraph, value)
            _shade_docx_element(paragraph._p, "F2F2F2")
        else:
            paragraph = document.add_paragraph()
            _add_docx_inline_text(paragraph, value)

        paragraph.paragraph_format.space_after = Pt(
            2 if opts["style"] == "Minimal"
            else (8 if opts["style"] == "Presentation" else 5)
        )
        paragraph.paragraph_format.line_spacing = (
            1.0 if opts["style"] == "Minimal"
            else (1.35 if opts["style"] == "Presentation" else 1.15)
        )
        previous_kind = kind

    _set_docx_cell_watermark(section, opts["watermark"])

    if opts["include_company_info"]:
        footer = section.footer.paragraphs[0]
        footer.text = f"{opts['company_name']} • {opts['company_info']}"
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer.runs[0].font.size = Pt(8)
        footer.runs[0].font.color.rgb = RGBColor(107, 114, 128)

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue(), max(1, len(document.paragraphs))


def _add_pptx_logo(slide, opts, left, top, width):
    if not opts["include_logo"]:
        return
    logo_stream = _prepared_logo_stream(opts["logo_path"])
    if logo_stream is None:
        raise RuntimeError(
            "AutoTecPro logo was selected, but App/logo.png could not be loaded."
        )
    slide.shapes.add_picture(logo_stream, left, top, width=width)


def _add_pptx_watermark(slide, text, presentation):
    if not text:
        return
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    shape = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(2.8),
        presentation.slide_width - Inches(2.4),
        Inches(1.0),
    )
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = str(text).upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(210, 210, 210)
    shape.rotation = 325


def _build_pptx_with_options(document_text, title, cleaner=None, options=None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    opts = _document_options(options)
    profile = _style_profile(opts["style"])
    text = _clean_document_text(document_text, cleaner)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    # No standalone cover/title slide. Content begins immediately.
    document_subtitle = title

    max_lines = (
        4 if opts["style"] == "Presentation"
        else (10 if opts["style"] == "Minimal" else 7)
    )
    body_font = (
        25 if opts["style"] == "Presentation"
        else (16 if opts["style"] == "Minimal" else 20)
    )

    for section_title, section_lines in _split_slide_sections(text):
        chunks = [
            section_lines[index:index + max_lines]
            for index in range(0, len(section_lines), max_lines)
        ] or [[]]

        for chunk_index, chunk in enumerate(chunks):
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[1]
            )
            if len(presentation.slides) == 1 and chunk_index == 0:
                slide_title = (
                    document_subtitle
                    if section_title == "Overview"
                    else f"{document_subtitle} — {section_title}"
                )
            else:
                slide_title = (
                    section_title
                    if chunk_index == 0
                    else f"{section_title} (continued)"
                )
            slide.shapes.title.text = slide_title
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(
                30 if opts["style"] == "Presentation"
                else (22 if opts["style"] == "Minimal" else 26)
            )
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = (
                RGBColor(185, 28, 28)
                if opts["style"] != "Minimal"
                else RGBColor(31, 41, 55)
            )

            frame = slide.placeholders[1].text_frame
            frame.clear()
            for item_index, line in enumerate(chunk):
                paragraph = (
                    frame.paragraphs[0]
                    if item_index == 0
                    else frame.add_paragraph()
                )
                paragraph.text = re.sub(
                    r"\*\*(.*?)\*\*",
                    r"\1",
                    line,
                )
                paragraph.font.size = Pt(body_font)
                paragraph.space_after = Pt(
                    10 if opts["style"] == "Presentation"
                    else (2 if opts["style"] == "Minimal" else 6)
                )

            _add_pptx_logo(
                slide,
                opts,
                Inches(11.05),
                Inches(0.13),
                Inches(1.55 if opts["style"] != "Minimal" else 1.1),
            )
            _add_pptx_watermark(slide, opts["watermark"], presentation)

            if opts["include_company_info"]:
                footer = slide.shapes.add_textbox(
                    Inches(0.45),
                    Inches(7.08),
                    Inches(12.2),
                    Inches(0.25),
                )
                p = footer.text_frame.paragraphs[0]
                p.text = f"{opts['company_name']} • {opts['company_info']}"
                p.font.size = Pt(7.5)
                p.font.color.rgb = RGBColor(107, 114, 128)
                p.alignment = PP_ALIGN.CENTER

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue(), len(presentation.slides)


def _build_xlsx_with_options(document_text, title, cleaner=None, options=None):
    from openpyxl import Workbook
    from openpyxl.styles import (
        Alignment, Font, PatternFill, Border, Side,
    )
    from openpyxl.utils import get_column_letter
    opts = _document_options(options)
    profile = _style_profile(opts["style"])
    text = _clean_document_text(document_text, cleaner)
    rows = _extract_table_rows(text)

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Document"

    width = max(1, max((len(row) for row in rows), default=1))
    sheet.merge_cells(
        start_row=1,
        start_column=1,
        end_row=1,
        end_column=width,
    )
    sheet["A1"] = title
    sheet["A1"].font = Font(
        bold=True,
        size=(
            22 if opts["style"] == "Presentation"
            else (14 if opts["style"] == "Minimal" else 17)
        ),
        color=(
            "B91C1C"
            if opts["style"] != "Minimal"
            else "1F2937"
        ),
    )
    sheet["A1"].alignment = Alignment(
        horizontal=(
            "center"
            if opts["style"] != "Minimal"
            else "left"
        )
    )
    sheet.row_dimensions[1].height = (
        34 if opts["style"] == "Presentation"
        else (22 if opts["style"] == "Minimal" else 28)
    )

    start_row = 3
    if opts["include_company_info"]:
        sheet.merge_cells(
            start_row=2,
            start_column=1,
            end_row=2,
            end_column=width,
        )
        sheet["A2"] = f"{opts['company_name']} • {opts['company_info']}"
        sheet["A2"].font = Font(size=9, color="6B7280")
        sheet["A2"].alignment = Alignment(horizontal="center")
        start_row = 4

    if opts["include_logo"]:
        logo_stream = _prepared_logo_stream(opts["logo_path"])
        if logo_stream is None:
            raise RuntimeError(
                "AutoTecPro logo was selected, but App/logo.png could not be loaded."
            )
        from openpyxl.drawing.image import Image as XLImage
        image = XLImage(logo_stream)
        image.width = (
            190 if opts["style"] == "Presentation"
            else (110 if opts["style"] == "Minimal" else 150)
        )
        image.height = int(image.width * 0.30)
        sheet.add_image(
            image,
            f"{get_column_letter(max(2, width))}1",
        )

    if opts["watermark"]:
        sheet.oddHeader.center.text = f"&B&KBBBBBB{opts['watermark'].upper()}"
        sheet.evenHeader.center.text = f"&B&KBBBBBB{opts['watermark'].upper()}"

    thin_border = Border(
        left=Side(style="thin", color="D1D5DB"),
        right=Side(style="thin", color="D1D5DB"),
        top=Side(style="thin", color="D1D5DB"),
        bottom=Side(style="thin", color="D1D5DB"),
    )

    for row_index, row in enumerate(rows, start=start_row):
        for column_index, value in enumerate(row, start=1):
            cell = sheet.cell(row_index, column_index, value=value)
            cell.alignment = Alignment(
                vertical="top",
                wrap_text=True,
            )

            if row_index == start_row and len(rows) > 1:
                cell.font = Font(
                    bold=True,
                    color="FFFFFF",
                    size=11 if opts["style"] != "Minimal" else 10,
                )
                cell.fill = PatternFill(
                    "solid",
                    fgColor=(
                        "B91C1C"
                        if opts["style"] != "Minimal"
                        else "374151"
                    ),
                )
            elif opts["style"] == "Presentation":
                cell.font = Font(size=12)
                cell.fill = PatternFill(
                    "solid",
                    fgColor=(
                        "F8FAFC"
                        if row_index % 2
                        else "EEF2F7"
                    ),
                )
            elif opts["style"] == "Minimal":
                cell.font = Font(size=9)
            else:
                cell.font = Font(size=10)
                cell.border = thin_border

    for column_index in range(1, width + 1):
        max_length = max(
            (
                len(str(sheet.cell(row, column_index).value or ""))
                for row in range(1, sheet.max_row + 1)
            ),
            default=10,
        )
        sheet.column_dimensions[
            get_column_letter(column_index)
        ].width = min(
            max(
                max_length + 2,
                11 if opts["style"] == "Minimal" else 13,
            ),
            48 if opts["style"] == "Presentation" else 40,
        )

    sheet.freeze_panes = f"A{start_row}"
    sheet.sheet_view.showGridLines = opts["style"] == "Minimal"

    if opts["include_company_info"]:
        sheet.oddFooter.center.text = (
            f"{opts['company_name']} • {opts['company_info']}"
        )

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue(), len(workbook.sheetnames)


def _build_csv_with_options(document_text, title, cleaner=None, options=None):
    opts = _document_options(options)
    text = _clean_document_text(document_text, cleaner)
    rows = _extract_table_rows(text)
    buffer = io.StringIO(newline="")
    writer = csv.writer(buffer)

    # CSV cannot contain visual styling, images, or a real watermark.
    # Include explicit text rows only when the relevant options are enabled.
    if opts["include_company_info"]:
        writer.writerow([opts["company_name"]])
        writer.writerow([opts["company_info"]])
        writer.writerow([])
    if opts["watermark"]:
        writer.writerow([opts["watermark"].upper()])
        writer.writerow([])

    writer.writerows(rows)
    return buffer.getvalue().encode("utf-8-sig"), max(1, len(rows))



def create_document_record(
    prompt_text: Any,
    answer_text: Any,
    format_name: str | None = None,
    visible_text_cleaner: Callable[[Any], str] | None = None,
    options: dict[str, Any] | None = None,
) -> dict[str, Any]:
    request = detect_document_generation_request(prompt_text)
    resolved_format = str(format_name or (request or {}).get("format") or "pdf").lower()
    if resolved_format not in FORMAT_CONFIG:
        raise ValueError(f"Unsupported document format: {resolved_format}")
    title = derive_document_title(prompt_text, answer_text)
    filename = safe_document_filename(title, answer_text, resolved_format)
    opts = _document_options(options)
    if resolved_format == "pdf": file_bytes, unit_count, unit_label = *_build_pdf_with_options(answer_text, title, visible_text_cleaner, opts), "Pages"
    elif resolved_format == "docx": file_bytes, unit_count, unit_label = *_build_docx_with_options(answer_text, title, visible_text_cleaner, opts), "Paragraphs"
    elif resolved_format == "pptx": file_bytes, unit_count, unit_label = *_build_pptx_with_options(answer_text, title, visible_text_cleaner, opts), "Slides"
    elif resolved_format == "xlsx": file_bytes, unit_count, unit_label = *_build_xlsx_with_options(answer_text, title, visible_text_cleaner, opts), "Sheets"
    else: file_bytes, unit_count, unit_label = *_build_csv_with_options(answer_text, title, visible_text_cleaner, opts), "Rows"
    config = FORMAT_CONFIG[resolved_format]; encoded = base64.b64encode(file_bytes).decode("ascii")
    return {
        "name": filename, "format": resolved_format, "format_label": config["label"],
        "icon": config["icon"], "mime_type": config["mime_type"],
        "data_url": f"data:{config['mime_type']};base64,{encoded}",
        "unit_count": int(unit_count), "unit_label": unit_label,
        "page_count": int(unit_count) if resolved_format == "pdf" else 0,
        "size_bytes": len(file_bytes), "generated": True,
        "style": opts["style"],
        "branded": bool(opts["include_logo"] or opts["include_company_info"]),
        "logo_requested": bool(opts["include_logo"]),
        "logo_available": bool(
            opts["logo_path"] and Path(str(opts["logo_path"])).is_file()
        ),
        "company_info_requested": bool(opts["include_company_info"]),
        "watermark_requested": bool(opts["watermark"]),
        "generator_version": "2.4",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
